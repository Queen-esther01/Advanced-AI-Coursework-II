import base64
import re
import shutil
import subprocess
import sys
from pathlib import Path, PurePosixPath

from docx import Document
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

NS_V = "urn:schemas-microsoft-com:vml"
VECTOR_IMAGE_EXTS = frozenset({"emf", "wmf"})


def _emf2svg_to_png(emf_path: Path, png_path: Path, errors: list[str]) -> Path | None:
    emf2svg = shutil.which("emf2svg-conv")
    rsvg = shutil.which("rsvg-convert")
    if not emf2svg or not rsvg:
        return None

    svg_path = emf_path.with_suffix(".svg")
    for label, extra in (("emf2svg+emfplus", ["-p"]), ("emf2svg", [])):
        try:
            subprocess.run(
                [emf2svg, "-i", str(emf_path), "-o", str(svg_path), *extra],
                check=True,
                capture_output=True,
            )
            if not svg_path.exists() or svg_path.stat().st_size == 0:
                errors.append(f"{label}: produced no SVG")
                continue
            subprocess.run(
                [rsvg, "-o", str(png_path), str(svg_path)],
                check=True,
                capture_output=True,
            )
            if png_path.exists() and png_path.stat().st_size > 0:
                return png_path
            errors.append(f"{label}: produced no PNG")
        except subprocess.CalledProcessError as exc:
            err = (exc.stderr or b"").decode(errors="replace").strip()
            errors.append(f"{label}: {err or exc}")
        finally:
            svg_path.unlink(missing_ok=True)
    return None


def _soffice_vector_to_png(
    emf_path: Path, png_path: Path, errors: list[str]
) -> Path | None:
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        return None
    out_dir = png_path.parent
    out_dir.mkdir(parents=True, exist_ok=True)
    try:
        result = subprocess.run(
            [
                soffice,
                "--headless",
                "--convert-to",
                "png",
                "--outdir",
                str(out_dir),
                str(emf_path.resolve()),
            ],
            capture_output=True,
            text=True,
        )
        if result.returncode != 0:
            errors.append(
                f"soffice: {(result.stderr or result.stdout or '').strip() or result.returncode}"
            )
            return None
        produced = out_dir / f"{emf_path.stem}.png"
        if not produced.is_file() or produced.stat().st_size == 0:
            errors.append(
                f"soffice: expected output missing ({produced})"
            )
            return None
        if produced.resolve() != png_path.resolve():
            shutil.copy2(produced, png_path)
        if png_path.exists() and png_path.stat().st_size > 0:
            return png_path
        errors.append("soffice: produced no PNG")
    except Exception as exc:
        errors.append(f"soffice: {exc}")
    return None


def convert_emf_to_png(
    emf_path: str | Path, png_path: str | Path | None = None
) -> Path:
    emf_path = Path(emf_path)
    if not emf_path.exists():
        raise FileNotFoundError(emf_path)
    png_path = Path(png_path) if png_path else emf_path.with_suffix(".png")
    png_path.parent.mkdir(parents=True, exist_ok=True)

    errors: list[str] = []
    result = _emf2svg_to_png(emf_path, png_path, errors)
    if result:
        return result

    result = _soffice_vector_to_png(emf_path, png_path, errors)
    if result:
        return result

    gs = shutil.which("gs")
    if gs:
        try:
            subprocess.run(
                [
                    gs,
                    "-dSAFER",
                    "-dBATCH",
                    "-dNOPAUSE",
                    "-sDEVICE=png16m",
                    "-r150",
                    f"-sOutputFile={png_path}",
                    str(emf_path),
                ],
                check=True,
                capture_output=True,
            )
            if png_path.exists() and png_path.stat().st_size > 0:
                return png_path
            errors.append("gs: produced no output")
        except subprocess.CalledProcessError as exc:
            err = (exc.stderr or b"").decode(errors="replace").strip()
            errors.append(f"gs: {err or exc}")

    magick = shutil.which("magick") or shutil.which("convert")
    if magick:
        try:
            subprocess.run(
                [magick, str(emf_path), str(png_path)],
                check=True,
                capture_output=True,
            )
            if png_path.exists() and png_path.stat().st_size > 0:
                return png_path
            errors.append("magick: produced no output")
        except subprocess.CalledProcessError as exc:
            err = (exc.stderr or b"").decode(errors="replace").strip()
            errors.append(f"magick: {err or exc}")

    raise RuntimeError(
        f"Could not convert {emf_path} to PNG. "
        f"{'; '.join(errors) or 'no converter found'}. "
        "Install: brew install libemf2svg librsvg"
    )


def _maybe_convert_vector_image(path: Path, warn: bool = True) -> Path:
    if path.suffix.lstrip(".").lower() not in VECTOR_IMAGE_EXTS:
        return path
    try:
        return convert_emf_to_png(path)
    except RuntimeError as exc:
        if warn:
            print(f"Warning: {exc}", file=sys.stderr)
        return path


def _escape_table_cell(text):
    return text.replace("|", "\\|").replace("\n", " ").strip()


def _clean_table_cell(text):
    return " ".join(str(text).replace("\n", " ").split()).strip()


def _is_empty_table_row(row: list[str]) -> bool:
    return all(not _clean_table_cell(c) for c in row)


def _is_separator_table_line(line: str) -> bool:
    stripped = line.strip()
    if not stripped.startswith("|"):
        return False
    cells = [_clean_table_cell(c) for c in stripped.strip("|").split("|")]
    return cells and all(
        re.fullmatch(r"-+", c.replace(" ", "")) or not c for c in cells
    )


def _join_cell_paragraphs(paragraphs: list[str]) -> str:
    parts = [" ".join(p.split()) for p in paragraphs if p and p.strip()]
    if not parts:
        return ""
    out = parts[0]
    for part in parts[1:]:
        if out.endswith((".", ":", ";")):
            sep = " "
        elif out[-1].isalpha() and part[0].isupper():
            sep = ". "
        elif out[-1].isalpha() and part[0].isdigit():
            sep = ". "
        else:
            sep = " "
        out = out + sep + part
    return out


def _table_cell_text(tc) -> str:
    paragraphs = []
    for para_el in tc.findall(qn("w:p")):
        text = "".join(t.text or "" for t in para_el.findall(f".//{qn('w:t')}"))
        if text.strip():
            paragraphs.append(text)
    if paragraphs:
        return _join_cell_paragraphs(paragraphs)
    return "".join(t.text or "" for t in tc.findall(f".//{qn('w:t')}")).strip()


def _table_cell_span(tc) -> tuple[int, str | None]:
    colspan = 1
    vmerge = None
    tc_pr = tc.find(qn("w:tcPr"))
    if tc_pr is not None:
        grid_span = tc_pr.find(qn("w:gridSpan"))
        if grid_span is not None:
            colspan = int(grid_span.get(qn("w:val"), 1))
        vm = tc_pr.find(qn("w:vMerge"))
        if vm is not None:
            vmerge = vm.get(qn("w:val"), "continue")
    return colspan, vmerge


def _extract_table_grid(tbl_el) -> list[list[str]]:
    grid: list[list[str]] = []
    vmerge_values: dict[int, str] = {}

    for row_el in tbl_el.findall(qn("w:tr")):
        row: list[str] = []
        col = 0
        for tc in row_el.findall(qn("w:tc")):
            text = _table_cell_text(tc)
            colspan, vmerge = _table_cell_span(tc)
            if vmerge == "continue":
                text = vmerge_values.get(col, "")
            else:
                vmerge_values[col] = text
            for offset in range(colspan):
                row.append(text if offset == 0 else text)
                col += 1
        grid.append(row)

    if not grid:
        return []

    width = max(len(row) for row in grid)
    for row in grid:
        while len(row) < width:
            row.append("")
    return grid


def _normalize_grid_rows(rows: list[list[str]]) -> list[list[str]]:
    if not rows:
        return []
    width = max(len(row) for row in rows)
    return [
        [_clean_table_cell(row[i]) if i < len(row) else "" for i in range(width)]
        for row in rows
    ]


def _is_probable_header_row(row: list[str], row_index: int) -> bool:
    cells = [_clean_table_cell(c) for c in row]
    if not any(cells):
        return False
    if row_index == 0:
        return True
    if any(c.upper() == "N/A" for c in cells):
        return False
    if max((len(c) for c in cells), default=0) >= 80:
        return False
    non_empty = [c for c in cells if c]
    return len(non_empty) >= 2 and all(len(c) < 80 for c in non_empty)


def _header_label(cell: str, index: int) -> str:
    cell = _clean_table_cell(cell)
    return cell or f"Column {index + 1}"


def _merge_header_rows(header_rows: list[list[str]]) -> list[str]:
    if not header_rows:
        return []
    if len(header_rows) == 1:
        return [_header_label(c, i) for i, c in enumerate(header_rows[0])]

    row0, row1 = header_rows[0], header_rows[1]
    row0_labels = [_clean_table_cell(c) for c in row0 if _clean_table_cell(c)]
    row0_generic = len(row0_labels) > 1 and len(set(row0_labels)) == 1

    width = max(len(r) for r in header_rows)
    merged = []
    for i in range(width):
        top = _clean_table_cell(row0[i]) if i < len(row0) else ""
        bottom = _clean_table_cell(row1[i]) if i < len(row1) else ""
        if not bottom:
            merged.append(_header_label(top, i))
        elif row0_generic or not top or top == bottom:
            merged.append(bottom)
        else:
            merged.append(f"{top} — {bottom}")
    return merged


def _count_header_rows(rows: list[list[str]]) -> int:
    count = 0
    for i, row in enumerate(rows):
        if _is_probable_header_row(row, i):
            count += 1
        else:
            break
    return max(count, 1) if rows else 0


def _first_column_is_row_label(headers: list[str]) -> bool:
    first = headers[0].lower() if headers else ""
    if not first:
        return True
    return "status" in first or "definition" in first


def _row_label(
    cells: list[str], carry_label: str | None
) -> tuple[str | None, str | None]:
    first = cells[0] if cells else ""
    if first and first.upper() != "N/A":
        return first, first
    return carry_label, carry_label


def _row_to_sentences(
    headers: list[str],
    row: list[str],
    carry_label: str | None,
    *,
    first_col_is_label: bool,
) -> tuple[list[str], str | None]:
    cells = [_clean_table_cell(c) for c in row]
    while len(cells) < len(headers):
        cells.append("")

    if _is_empty_table_row(cells):
        return [], carry_label

    if first_col_is_label:
        label, carry_label = _row_label(cells, carry_label)
    else:
        label = None
    start_col = 1 if first_col_is_label else 0
    groups: list[tuple[list[str], str]] = []

    for i in range(start_col, len(headers)):
        header = headers[i] if i < len(headers) else f"Column {i + 1}"
        cell = cells[i] if i < len(cells) else ""
        if not header or not cell:
            continue
        if first_col_is_label and groups and groups[-1][1] == cell:
            groups[-1][0].append(header)
        else:
            groups.append(([header], cell))

    sentences = []
    for group_headers, value in groups:
        combined_header = " | ".join(group_headers)
        if first_col_is_label and label:
            sentences.append(f"{label} — {combined_header}: {value}")
        else:
            sentences.append(f"{combined_header}: {value}")

    return sentences, carry_label


def _table_rows_to_sentences(rows: list[list[str]]) -> str:
    if not rows:
        return ""
    cleaned = _normalize_grid_rows(rows)
    cleaned = [row for row in cleaned if not _is_empty_table_row(row)]
    if not cleaned:
        return ""

    header_count = _count_header_rows(cleaned)
    headers = _merge_header_rows(cleaned[:header_count])
    first_col_is_label = _first_column_is_row_label(headers)
    sentences: list[str] = []
    carry_label: str | None = None

    for row in cleaned[header_count:]:
        row_sentences, carry_label = _row_to_sentences(
            headers, row, carry_label, first_col_is_label=first_col_is_label
        )
        sentences.extend(row_sentences)

    return "\n".join(sentences)


def _parse_markdown_table_row(line: str) -> list[str]:
    return [_clean_table_cell(c) for c in line.strip().strip("|").split("|")]


def deserialize_tables_in_markdown(md: str) -> str:
    lines = md.splitlines()
    out: list[str] = []
    i = 0
    while i < len(lines):
        line = lines[i]
        if (
            i + 1 < len(lines)
            and line.strip().startswith("|")
            and _is_separator_table_line(lines[i + 1])
        ):
            table_rows = [_parse_markdown_table_row(line)]
            i += 2
            while i < len(lines) and lines[i].strip().startswith("|"):
                table_rows.append(_parse_markdown_table_row(lines[i]))
                i += 1
            block = _table_rows_to_sentences(table_rows)
            if block:
                out.extend(block.splitlines())
            continue
        out.append(line)
        i += 1
    return "\n".join(out)


def _table_to_markdown(rows):
    if not rows:
        return ""
    lines = []
    header = rows[0]
    lines.append("| " + " | ".join(_escape_table_cell(c) for c in header) + " |")
    lines.append("| " + " | ".join("---" for _ in header) + " |")
    for row in rows[1:]:
        lines.append("| " + " | ".join(_escape_table_cell(c) for c in row) + " |")
    return "\n".join(lines)


def _heading_level(style_name):
    if style_name.startswith("Heading"):
        try:
            return int(style_name.split()[-1])
        except ValueError:
            return 1
    if style_name in ("Sub header", "Subheader", "TOC Heading"):
        return 2
    return None


def _paragraph_to_markdown(style, text):
    text = text.strip()
    if not text:
        return ""
    level = _heading_level(style)
    if level:
        return f"{'#' * level} {text}"
    return text


def _image_to_markdown(path):
    path = Path(path)
    return f"![{path.name}](<{path.resolve()}>)"


IMAGE_MARKDOWN = re.compile(r"^!\[(?P<alt>.+?)\]\(<(?P<path>.+?)>\)\s*$")


def parse_image_markdown(line: str) -> tuple[str, Path] | None:
    match = IMAGE_MARKDOWN.match(line.strip())
    if not match:
        return None
    return match.group("alt"), Path(match.group("path"))


def _libreoffice_convert(source: Path, out_dir: Path, fmt: str) -> Path:
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        raise RuntimeError(
            "LibreOffice not found. Install it (e.g. brew install --cask libreoffice) "
            "to read legacy Word .doc files."
        )
    out_dir.mkdir(parents=True, exist_ok=True)
    result = subprocess.run(
        [
            soffice,
            "--headless",
            "--convert-to",
            fmt,
            "--outdir",
            str(out_dir),
            str(source.resolve()),
        ],
        capture_output=True,
        text=True,
    )
    if result.returncode != 0:
        raise RuntimeError(f"LibreOffice failed:\n{result.stderr or result.stdout}")
    ext = f".{fmt}" if not fmt.startswith(".") else fmt
    converted = out_dir / f"{source.stem}{ext}"
    if not converted.is_file():
        raise FileNotFoundError(f"Expected converted file not found: {converted}")
    return converted


def _image_output_dir(source, image_dir=None):
    source = Path(source)
    if image_dir is not None:
        return Path(image_dir)
    return source.parent / f"{source.stem}_media"


def _output_path(directory, name):
    return Path(directory) / name


def _safe_filename(name):
    return "".join(c if c.isalnum() or c in "._-" else "_" for c in name)


def _image_rids_from_element(elem):
    rids = []
    seen = set()
    for blip in elem.findall(f".//{qn('a:blip')}"):
        rid = blip.get(qn("r:embed"))
        if rid and rid not in seen:
            seen.add(rid)
            rids.append(rid)
    for imagedata in elem.findall(f".//{{{NS_V}}}imagedata"):
        rid = imagedata.get(qn("r:id"))
        if rid and rid not in seen:
            seen.add(rid)
            rids.append(rid)
    return rids


def _paragraph_style(elem):
    p_pr = elem.find(qn("w:pPr"))
    if p_pr is None:
        return "Normal"
    p_style = p_pr.find(qn("w:pStyle"))
    if p_style is None:
        return "Normal"
    return p_style.get(qn("w:val"), "Normal")


def _extract_docx_images(doc, output_dir, convert_vector=True):
    output_dir.mkdir(parents=True, exist_ok=True)
    paths = {}
    for rel_id, rel in doc.part.rels.items():
        if "image" not in rel.reltype:
            continue
        name = PurePosixPath(rel.target_ref).name
        dest = _output_path(output_dir, name)
        dest.write_bytes(rel.target_part.blob)
        if convert_vector:
            dest = _maybe_convert_vector_image(dest)
        paths[rel_id] = dest.resolve()
    return paths


def _join_blocks(blocks):
    return "\n\n".join(b for b in blocks if b)


def _blocks_to_markdown(blocks):
    parts = []
    for block in blocks:
        if block["type"] == "image" and block.get("saved_to"):
            parts.append(_image_to_markdown(block["saved_to"]))
        elif block["type"] == "text":
            md = _paragraph_to_markdown(block["style"] or "Normal", block["content"])
            if md:
                parts.append(md)
    return _join_blocks(parts)


def extract_blocks(docx_path, images_dir=None, convert_vector=True):
    source = Path(docx_path)
    doc = Document(source)
    _images_dir = Path(images_dir) if images_dir else _image_output_dir(source)
    image_paths = _extract_docx_images(doc, _images_dir, convert_vector=convert_vector)

    blocks = []
    idx = 0

    def add_image(rid):
        nonlocal idx
        rel = doc.part.rels.get(rid)
        if rel is None or "image" not in rel.reltype or rid not in image_paths:
            return
        path = Path(image_paths[rid])
        blocks.append(
            {
                "index": idx,
                "type": "image",
                "style": None,
                "content": base64.b64encode(path.read_bytes()).decode("utf-8"),
                "encoding": "base64",
                "filename": path.name,
                "ext": path.suffix.lstrip(".").lower(),
                "saved_to": str(path),
            }
        )
        idx += 1

    def add_text(text, style):
        nonlocal idx
        text = text.strip()
        if not text:
            return
        blocks.append(
            {
                "index": idx,
                "type": "text",
                "style": style,
                "content": text,
                "encoding": "utf-8",
                "filename": None,
                "ext": None,
                "saved_to": None,
            }
        )
        idx += 1

    for child in doc.element.body.iterchildren():
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag

        if tag == "p":
            for rid in _image_rids_from_element(child):
                add_image(rid)
            text = "".join(t.text or "" for t in child.findall(f".//{qn('w:t')}"))
            add_text(text, _paragraph_style(child))
        elif tag == "tbl":
            for cell_el in child.findall(f".//{qn('w:tc')}"):
                for para_el in cell_el.findall(f".//{qn('w:p')}"):
                    for rid in _image_rids_from_element(para_el):
                        add_image(rid)
            grid = _extract_table_grid(child)
            add_text(_table_rows_to_sentences(grid), "Table")

    return blocks


def read_docx(file, image_dir=None, convert_vector=True):
    source = Path(file)
    images_dir = _image_output_dir(source, image_dir)
    blocks = extract_blocks(
        source, images_dir=images_dir, convert_vector=convert_vector
    )
    return _blocks_to_markdown(blocks)


def read_doc(file, image_dir=None, convert_vector=True):
    source = Path(file)
    convert_dir = source.parent / f"{source.stem}_converted"
    docx_path = convert_dir / f"{source.stem}.docx"
    if not docx_path.is_file() or docx_path.stat().st_mtime < source.stat().st_mtime:
        _libreoffice_convert(source, convert_dir, "docx")
    return read_docx(
        docx_path,
        image_dir=image_dir or _image_output_dir(source),
        convert_vector=convert_vector,
    )


_RASTER_IMAGE_EXTS = frozenset({"png", "jpg", "jpeg", "gif", "webp"})


def _extract_pptx_image(shape, output_dir, convert_vector: bool = True):
    image = shape.image
    ext = (image.ext or "bin").lower()
    filename = f"{_safe_filename(shape.name or 'image')}_{shape.shape_id}.{ext}"
    dest = _output_path(output_dir, filename)
    dest.write_bytes(image.blob)
    dest = dest.resolve()

    png_path = dest.with_suffix(".png")
    if png_path.is_file() and png_path.stat().st_size > 0:
        return png_path.resolve()

    if ext in VECTOR_IMAGE_EXTS:
        if not convert_vector:
            return None
        converted = _maybe_convert_vector_image(dest, warn=True)
        if converted.suffix.lstrip(".").lower() in _RASTER_IMAGE_EXTS:
            return converted.resolve()
        return None

    if ext in _RASTER_IMAGE_EXTS:
        return dest

    return None


def _is_picture_shape(shape):
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return True
    try:
        shape.image
        return True
    except Exception:
        return False


def _pptx_shapes_to_markdown(shapes, output_dir, convert_vector: bool = True):
    elements = []
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            elements.extend(
                _pptx_shapes_to_markdown(shape.shapes, output_dir, convert_vector)
            )
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table_data = [[cell.text for cell in row.cells] for row in shape.table.rows]
            text = _table_rows_to_sentences(table_data)
            if text:
                elements.append(text)
        elif _is_picture_shape(shape):
            path = _extract_pptx_image(shape, output_dir, convert_vector)
            if path:
                elements.append(_image_to_markdown(path))
        elif hasattr(shape, "text") and shape.text.strip():
            elements.append(shape.text.strip())
        elif shape.has_chart:
            elements.append(f"[chart: {shape.chart.chart_type}]")
    return elements


def read_pptx(file, image_dir=None, convert_vector: bool = True):
    source = Path(file)
    output_dir = _image_output_dir(source, image_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    prs = Presentation(source)
    slides = []

    for slide_index, slide in enumerate(prs.slides):
        elements = _pptx_shapes_to_markdown(
            slide.shapes, output_dir, convert_vector=convert_vector
        )
        slide_body = _join_blocks(elements)
        if slide_body:
            slides.append(f"## Slide {slide_index + 1}\n\n{slide_body}")

    return _join_blocks(slides)


def read_pptx_slides(
    file: Path | str,
    slides_dir: Path | str | None = None,
    *,
    dpi: int = 150,
) -> tuple[str, list[dict]]:
    """
    Render each slide as one JPEG and return markdown (slide text only) plus a manifest.
    Use this for CPT decks where the slide composition is the meaningful unit.
    """
    from .slide_renderer import manifest_to_markdown, render_slides

    source = Path(file)
    if slides_dir is None:
        slides_dir = source.parent / f"{source.stem}_slides"
    manifest = render_slides(source, Path(slides_dir), dpi=dpi)
    return manifest_to_markdown(manifest), manifest


def chunk_markdown(md: str, metadata: dict):
    max_chars = 1500
    overlap_chars = 150

    def _split_large_section(text: str) -> list[str]:
        if len(text) <= max_chars:
            return [text]

        lines = text.splitlines()
        if not lines:
            return [text]

        heading = lines[0]
        body = "\n".join(lines[1:]).strip()
        if not body:
            return [text]

        paragraphs = [p.strip() for p in re.split(r"\n\s*\n", body) if p.strip()]
        if not paragraphs:
            return [text]

        windows: list[str] = []
        current = heading

        for paragraph in paragraphs:
            candidate = f"{current}\n\n{paragraph}" if current else paragraph
            if len(candidate) <= max_chars:
                current = candidate
                continue

            if current and current != heading:
                windows.append(current)
            elif len(candidate) > max_chars:
                start = 0
                while start < len(paragraph):
                    end = min(start + max_chars, len(paragraph))
                    piece = paragraph[start:end]
                    window = f"{heading}\n\n{piece}"
                    windows.append(window)
                    if end == len(paragraph):
                        break
                    start = max(0, end - overlap_chars)
                current = heading
                continue

            current = f"{heading}\n\n{paragraph}"

        if current and current != heading:
            windows.append(current)

        return windows or [text]

    md = deserialize_tables_in_markdown(md.strip())
    parts = re.split(r"(?=(?:(?<=\n)|^)#{1,6} )", md)
    chunks = []
    for part in parts:
        part = part.strip()
        if not part or re.match(r"^#{1,6}\s+Contents\b", part):
            continue
        title = re.sub(r"^#{1,6}\s+", "", part.split("\n", 1)[0]).strip()
        split_parts = _split_large_section(part)
        if len(split_parts) == 1:
            chunks.append(
                {
                    "text": split_parts[0],
                    "metadata": {**metadata, "section": title},
                }
            )
            continue

        for i, split_part in enumerate(split_parts, start=1):
            chunks.append(
                {
                    "text": split_part,
                    "metadata": {**metadata, "section": title, "part": i},
                }
            )
    return chunks


def chunks_for_chroma(chunks):
    return [c["text"] for c in chunks], [c["metadata"] for c in chunks]
